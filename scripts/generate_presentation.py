#!/usr/bin/env python3
"""
AI-Wireshark Analyzer — Presentation Generator
Generates a professional PPTX presentation covering tool importance, features,
usage, advantages, and roadmap.

Usage:
    python3 scripts/generate_presentation.py [output_path]

Output defaults to: docs/AI-Wireshark-Analyzer-Presentation.pptx
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ─────────────────────────────────────────────────────────────────────────────
# Colour palette
# ─────────────────────────────────────────────────────────────────────────────
C_NAVY      = RGBColor(0x1A, 0x25, 0x2F)   # slide backgrounds (dark navy)
C_BLUE      = RGBColor(0x2C, 0x3E, 0x50)   # secondary dark
C_ACCENT    = RGBColor(0x3B, 0x82, 0xF6)   # blue accent
C_GREEN     = RGBColor(0x22, 0xC5, 0x5E)   # success / positive
C_ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)   # warning / highlight
C_RED       = RGBColor(0xEF, 0x44, 0x44)   # critical
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xE2, 0xE8, 0xF0)   # light text / borders
C_SUBTLE    = RGBColor(0x94, 0xA3, 0xB8)   # subtle text
C_SLIDE_BG  = RGBColor(0xF4, 0xF6, 0xF9)   # light slide background

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# Helper utilities
# ─────────────────────────────────────────────────────────────────────────────

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=C_WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_text(slide, lines, left, top, width, height,
                    font_size=14, color=C_LIGHT, title=None, title_color=C_WHITE,
                    title_size=16, bullet_char="▸ ", spacing_before=6):
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    para_idx = 0
    if title:
        p = tf.paragraphs[para_idx]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
        para_idx += 1

    for i, line in enumerate(lines):
        if para_idx == 0 and i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if spacing_before:
            p.space_before = Pt(spacing_before)
        run = p.add_run()
        run.text = f"{bullet_char}{line}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

    return txBox


def set_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_divider_line(slide, left, top, width, color=C_ACCENT, thickness=3):
    from pptx.util import Pt
    line = slide.shapes.add_connector(1, Inches(left), Inches(top),
                                       Inches(left + width), Inches(top))
    line.line.color.rgb = color
    line.line.width = Pt(thickness)
    return line


def add_card(slide, left, top, width, height, title, body_lines,
             card_bg=C_BLUE, title_color=C_ACCENT, body_color=C_LIGHT,
             title_size=13, body_size=11, icon=""):
    """Add a dark card with title and bullet lines."""
    add_shape(slide, left, top, width, height, fill_color=card_bg)
    # Title bar accent strip
    add_shape(slide, left, top, 0.05, height, fill_color=C_ACCENT)

    title_text = f"{icon} {title}" if icon else title
    add_text_box(slide, title_text, left + 0.12, top + 0.08,
                 width - 0.18, 0.35,
                 font_size=title_size, bold=True, color=title_color)

    add_bullet_text(slide, body_lines, left + 0.12, top + 0.42,
                    width - 0.22, height - 0.55,
                    font_size=body_size, color=body_color, bullet_char="• ",
                    spacing_before=3)


def add_stat_box(slide, left, top, value, label, color=C_ACCENT,
                 val_size=28, lbl_size=11):
    """Big number + label stat box."""
    add_shape(slide, left, top, 2.2, 1.2, fill_color=C_BLUE)
    add_shape(slide, left, top, 2.2, 0.06, fill_color=color)
    add_text_box(slide, value, left + 0.1, top + 0.1, 2.0, 0.7,
                 font_size=val_size, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, left + 0.1, top + 0.82, 2.0, 0.35,
                 font_size=lbl_size, color=C_SUBTLE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 — Title & Identity"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_background(slide, C_NAVY)

    # Left accent bar
    add_shape(slide, 0, 0, 0.12, 7.5, fill_color=C_ACCENT)

    # Right decorative panel
    add_shape(slide, 9.8, 0, 3.53, 7.5, fill_color=RGBColor(0x12, 0x1A, 0x22))

    # Icon area (large letter W)
    add_text_box(slide, "📡", 10.2, 1.5, 2.5, 2.0,
                 font_size=72, bold=False, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Main title
    add_text_box(slide, "AI-Wireshark Analyzer",
                 0.4, 1.6, 9.0, 1.4,
                 font_size=44, bold=True, color=C_WHITE)

    # Subtitle
    add_text_box(slide,
                 "Intelligent Network Capture Analysis — Wi-Fi · TCP/UDP · DNS · IPv6 · ML Anomaly",
                 0.4, 3.1, 9.0, 0.7,
                 font_size=18, color=C_ACCENT)

    # Version badge
    add_shape(slide, 0.4, 3.9, 2.2, 0.45, fill_color=C_ACCENT)
    add_text_box(slide, "v1.7.1  |  July 2026",
                 0.45, 3.91, 2.1, 0.42,
                 font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Tag line
    add_text_box(slide, "All analysis is local · No data leaves the machine · MIT License",
                 0.4, 6.7, 9.0, 0.5,
                 font_size=11, color=C_SUBTLE)


def slide_problem(prs):
    """Slide 2 — The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)
    add_divider_line(slide, 0.5, 1.45, 12.33)

    add_text_box(slide, "The Problem", 0.5, 0.25, 12.0, 0.8,
                 font_size=34, bold=True, color=C_WHITE)
    add_text_box(slide, "Network packet captures are complex — existing tools require deep expertise",
                 0.5, 1.05, 12.0, 0.5, font_size=16, color=C_SUBTLE)

    problems = [
        ("📋  Manual Analysis Burden",
         ["Wireshark requires expert knowledge to interpret captures",
          "Hundreds of 802.11 status/reason codes with no contextual help",
          "Remediations must be looked up manually in IEEE 802.11 specs"]),
        ("🔍  No Root Cause Clarity",
         ["\"Wrong password\" vs \"AP state-machine bug\" look identical in raw frames",
          "WPA3-SAE, 4-way EAPOL, PMKID caching failures are hard to distinguish",
          "TCP retransmissions from congestion vs. MTU black-hole are indistinguishable"]),
        ("🌐  Protocol Complexity",
         ["Wi-Fi, IPv6, TCP, DNS, DHCP each need specialist tools",
          "Cross-protocol correlation (e.g., DHCP failure after 802.11 success) is manual",
          "Security threats (SYN flood, DNS tunneling) hidden in high packet volumes"]),
        ("⏱️  Time & Cost",
         ["Hours spent on each incident without automation",
          "On-site specialist required for enterprise Wi-Fi troubleshooting",
          "No consistent reproducible report format across team members"]),
    ]

    cols = [(0.4, 2.0), (6.7, 2.0)]
    rows = [1.6, 4.0]
    idx = 0
    for row in rows:
        for col_l, col_w in cols:
            if idx < len(problems):
                title, bullets = problems[idx]
                add_card(slide, col_l, row, 5.9, 2.2, title, bullets,
                         title_size=13, body_size=11)
                idx += 1


def slide_solution(prs):
    """Slide 3 — The Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)
    add_divider_line(slide, 0.5, 1.45, 12.33)

    add_text_box(slide, "The Solution", 0.5, 0.25, 12.0, 0.8,
                 font_size=34, bold=True, color=C_WHITE)
    add_text_box(slide, "AI-Wireshark Analyzer — a standalone desktop application for capture analysis",
                 0.5, 1.05, 12.0, 0.5, font_size=16, color=C_SUBTLE)

    # Central value proposition
    add_shape(slide, 0.4, 1.65, 12.53, 1.3, fill_color=RGBColor(0x1E, 0x40, 0xAF))
    add_shape(slide, 0.4, 1.65, 0.08, 1.3, fill_color=C_ACCENT)
    add_text_box(slide,
                 "Open a PCAP/PCAPNG file → Get root-cause analysis with specific IPs, rates, IEEE spec references, and targeted fix steps — in seconds",
                 0.6, 1.75, 12.0, 1.0,
                 font_size=16, bold=True, color=C_WHITE)

    pillars = [
        ("🧠  Intelligent Analysis",
         "Moves beyond raw packet display to root-cause diagnosis using frame-sequence evidence, EAPOL handshake analysis, and IEEE 802.11-2020 spec references"),
        ("🔒  Fully Local & Private",
         "No cloud, no telemetry, no network calls during analysis. Captures never leave the analyst's machine — critical for confidential enterprise Wi-Fi data"),
        ("📊  Rich Actionable Reports",
         "HTML + JSON output with capture-specific IPs, rates, and numbered remediation steps. Connection failures always surface before RF statistics"),
        ("🛠️  Zero-Config Setup",
         "Single binary for Linux (1.1 GB ZIP). Requires only tshark. Works on any PCAP from Wireshark, tcpdump, or live capture scripts"),
    ]

    top = 3.15
    for i, (title, body) in enumerate(pillars):
        left = 0.4 + (i % 2) * 6.47
        t = top if i < 2 else top + 1.75
        add_shape(slide, left, t, 5.9, 1.55, fill_color=C_BLUE)
        add_shape(slide, left, t, 0.06, 1.55, fill_color=C_ACCENT)
        add_text_box(slide, title, left + 0.15, t + 0.1, 5.5, 0.45,
                     font_size=13, bold=True, color=C_ACCENT)
        add_text_box(slide, body, left + 0.15, t + 0.55, 5.6, 0.92,
                     font_size=11, color=C_LIGHT, wrap=True)


def slide_architecture(prs):
    """Slide 4 — Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)
    add_divider_line(slide, 0.5, 1.45, 12.33)

    add_text_box(slide, "Architecture", 0.5, 0.25, 12.0, 0.8,
                 font_size=34, bold=True, color=C_WHITE)
    add_text_box(slide, "Modular Python 3.12 platform — two ingestion paths, local-only processing",
                 0.5, 1.05, 12.0, 0.5, font_size=16, color=C_SUBTLE)

    # Architecture flow diagram (text-based)
    # Input box
    add_shape(slide, 4.4, 1.65, 4.5, 0.6, fill_color=C_ACCENT)
    add_text_box(slide, "📁  PCAP / PCAPNG Input", 4.45, 1.68, 4.4, 0.55,
                 font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Arrow down
    add_text_box(slide, "↙             ↘", 3.8, 2.35, 5.8, 0.45,
                 font_size=20, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Two path boxes
    add_shape(slide, 0.4, 2.8, 5.5, 0.8, fill_color=C_BLUE)
    add_text_box(slide, "PyShark Path\n(ML pipeline, protocol analyzers)",
                 0.5, 2.82, 5.3, 0.76, font_size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

    add_shape(slide, 7.4, 2.8, 5.5, 0.8, fill_color=C_BLUE)
    add_text_box(slide, "tshark-native Path\n(GUI panels, scripts — high performance)",
                 7.5, 2.82, 5.3, 0.76, font_size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Analyzer boxes
    analyzers = [
        (0.4, "🔬  Protocol Analyzers\nTCP · UDP · DNS · ICMP · DHCP · WLAN"),
        (3.5, "🤖  ML Models\nIsolation Forest · Autoencoder"),
        (6.6, "📡  WLAN Analyzer\n802.11 Auth/Assoc/SAE/RF"),
        (9.7, "📈  IPv6 · TCP/UDP Reports\nChannel Monitor · Combined"),
    ]
    for left, text in analyzers:
        add_shape(slide, left, 3.85, 2.8, 0.9, fill_color=RGBColor(0x12, 0x2A, 0x44))
        add_shape(slide, left, 3.85, 2.8, 0.07, fill_color=C_GREEN)
        add_text_box(slide, text, left + 0.1, 3.9, 2.6, 0.82,
                     font_size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Arrow down
    add_text_box(slide, "↓", 6.3, 4.85, 1.0, 0.4,
                 font_size=22, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Report engine
    add_shape(slide, 2.0, 5.3, 9.3, 0.7, fill_color=RGBColor(0x1E, 0x40, 0xAF))
    add_text_box(slide, "📋  HTMLReportGenerator — Dynamic RCA · Threat Ordering · Per-IP Remediation",
                 2.1, 5.34, 9.1, 0.62, font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Output
    outputs = [("🖥️  Desktop GUI\nPyQt6 — 9 panels", 1.0),
               ("⌨️  CLI\nClick + Rich", 4.5),
               ("🌐  REST API\nFastAPI", 7.5),
               ("📂  results/\nHTML + JSON", 10.2)]
    for label, left in outputs:
        add_shape(slide, left, 6.2, 2.2, 0.9, fill_color=C_BLUE)
        add_text_box(slide, label, left + 0.1, 6.22, 2.0, 0.86,
                     font_size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)


def slide_wlan(prs):
    """Slide 5 — WLAN Analysis Deep-Dive"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)
    add_divider_line(slide, 0.5, 1.45, 12.33)

    add_text_box(slide, "WLAN Analysis", 0.5, 0.25, 12.0, 0.8,
                 font_size=34, bold=True, color=C_WHITE)
    add_text_box(slide, "Deepest open-source 802.11 diagnostic engine available",
                 0.5, 1.05, 12.0, 0.5, font_size=16, color=C_SUBTLE)

    features = [
        ("🔐  Authentication & Association",
         ["40+ IEEE 802.11-2020 status codes with root-cause labels",
          "WPA3-SAE: Commit/Confirm phase tracking, anti-clogging token detection",
          "PMKID cache stale detection (Status 53 + Open auth pattern)",
          "4-way EAPOL Msg1→Msg4 handshake stall detection (wrong PSK)"]),
        ("📡  RF Performance",
         ["High retry rate with per-AP (BSSID) and per-channel breakdown",
          "RTS/CTS ratio analysis — hidden node indicator",
          "Signal strength trends (dBm), noise floor, SNR margin",
          "Power-save null-frame storm and PS transition analysis"]),
        ("🔗  Connection Lifecycle",
         ["Per-client session timelines (scan → auth → assoc → EAPOL → data)",
          "Connection delay root cause: multi-band scanning, weak signal hesitation",
          "Channel-switch overhead, passive scan dwell time",
          "IP connectivity failure detection after successful WPA2 handshake"]),
        ("🛡️  Security & WPA3",
         ["WPA3-SAE failure sessions with IEEE 802.11-2020 spec citations",
          "Stale-association deadlock (Status 30 + CCMP PN forensic evidence)",
          "Beacon loss, probe response imbalance, unprotected data frame detection",
          "Only surfaces actual failures — advisory notices suppressed"]),
    ]

    cols = [(0.4, 5.9), (6.8, 5.9)]
    rows = [1.6, 4.0]
    idx = 0
    for row in rows:
        for col_l, _ in cols:
            if idx < len(features):
                title, bullets = features[idx]
                add_card(slide, col_l, row, 5.9, 2.25, title, bullets,
                         title_size=13, body_size=11)
                idx += 1


def slide_protocol(prs):
    """Slide 6 — Protocol Analyzers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)
    add_divider_line(slide, 0.5, 1.45, 12.33)

    add_text_box(slide, "Protocol Analyzers", 0.5, 0.25, 12.0, 0.8,
                 font_size=34, bold=True, color=C_WHITE)
    add_text_box(slide,
                 "TCP · UDP · DNS · ICMP · DHCP — with Dynamic Root Cause Analysis (v1.7.1)",
                 0.5, 1.05, 12.0, 0.5, font_size=16, color=C_SUBTLE)

    protocols = [
        ("🔵  TCP", C_ACCENT,
         ["SYN flood: rate + top targets + rate-limit commands",
          "Port scan: scanner IPs + ports probed",
          "Retransmissions: affected path IPs + ethtool commands",
          "Zero window: slow consumer IPs + buffer tuning",
          "RST storm: affected ports + refused connections",
          "HTTP/HTTPS via TCP port filter (80, 443, 8080…)"]),
        ("🟢  UDP / DHCP", C_GREEN,
         ["UDP flood: top sources + targets + iptables commands",
          "Amplification: abused services (DNS/NTP) + victim IPs",
          "DHCP starvation: unique MAC count + snooping commands",
          "Rogue DHCP: unauthorized server IPs listed",
          "Fragmentation attack: variance + tiny fragment count",
          "UDP port scanning with scanner breakdown"]),
        ("🟡  DNS", C_ORANGE,
         ["Tunneling: high-entropy domain samples",
          "DGA / malware C2: algorithmically-generated name detection",
          "NXDOMAIN excess: top failing domains listed",
          "DNS amplification: source IPs + RRL configuration",
          "Cache poisoning: conflicting-response detection",
          "Sample queries always shown for analyst review"]),
        ("🔴  ICMP / Security", C_RED,
         ["ICMP flood + ping sweep: source/target IPs",
          "Ping of Death: oversized packets + attacker IPs",
          "Smurf attack: victim IP + amplifier count",
          "ICMP tunneling: payload anomaly + endpoint IPs",
          "SQL injection / XSS / path traversal: attacker IPs + URIs",
          "Suspicious user agents: scanner tool identification"]),
    ]

    left_positions = [0.4, 3.55, 6.7, 9.85]
    for i, (title, color, bullets) in enumerate(protocols):
        left = left_positions[i]
        add_shape(slide, left, 1.6, 3.0, 5.55, fill_color=C_BLUE)
        add_shape(slide, left, 1.6, 0.06, 5.55, fill_color=color)
        add_text_box(slide, title, left + 0.15, 1.68, 2.7, 0.42,
                     font_size=14, bold=True, color=color)
        add_bullet_text(slide, bullets, left + 0.15, 2.15,
                        2.7, 4.85, font_size=10.5, color=C_LIGHT,
                        bullet_char="→ ", spacing_before=4)


def slide_reporting(prs):
    """Slide 7 — Reporting & Threat Intelligence"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)
    add_divider_line(slide, 0.5, 1.45, 12.33)

    add_text_box(slide, "Intelligent Reporting", 0.5, 0.25, 12.0, 0.8,
                 font_size=34, bold=True, color=C_WHITE)
    add_text_box(slide, "Dynamic root cause + specific fix steps, not generic boilerplate",
                 0.5, 1.05, 12.0, 0.5, font_size=16, color=C_SUBTLE)

    # Left panel: features
    features = [
        ("📊  Threat Overview Table",
         ["Connection failures always above RF statistics",
          "Smart ordering: (severity, type) — critical/connection → medium/RF",
          "INFO-only advisory notices automatically suppressed",
          "Consistent across WLAN, TCP, UDP, DNS, ICMP, DHCP"]),
        ("🔬  Per-Threat Root Cause Analysis",
         ["25+ threat types have capture-specific RCA",
          "Actual IPs, rates, and counts — not generic placeholder text",
          "Specific fix commands with actual parameters (iptables, sysctl, tshark)",
          "Static guide used only as fallback when no dynamic data available"]),
    ]

    for i, (title, bullets) in enumerate(features):
        add_card(slide, 0.4, 1.65 + i * 2.55, 6.3, 2.35, title, bullets,
                 title_size=13, body_size=11)

    # Right panel: example RCA
    add_shape(slide, 7.0, 1.65, 6.0, 5.55, fill_color=RGBColor(0x0A, 0x14, 0x1E))
    add_shape(slide, 7.0, 1.65, 6.0, 0.06, fill_color=C_ACCENT)
    add_text_box(slide, "Example: SYN Flood RCA", 7.1, 1.72, 5.7, 0.42,
                 font_size=13, bold=True, color=C_ACCENT)

    # RCA box (blue)
    add_shape(slide, 7.1, 2.22, 5.8, 1.8, fill_color=RGBColor(0x1E, 0x3A, 0x5F))
    add_shape(slide, 7.1, 2.22, 0.06, 1.8, fill_color=RGBColor(0x3B, 0x82, 0xF6))
    add_text_box(slide, "🔍 Root Cause Analysis: SYN Flood",
                 7.2, 2.27, 5.6, 0.4, font_size=11, bold=True, color=C_ACCENT)
    rca_lines = ["• SYN rate: 1,500 pkts/sec (threshold: 100)",
                 "• Targeted servers: 10.0.1.45 (1,200)  10.0.1.46 (300)"]
    add_bullet_text(slide, rca_lines, 7.2, 2.68, 5.6, 1.2,
                    font_size=10, color=C_LIGHT, bullet_char="", spacing_before=5)

    # Recs box (green)
    add_shape(slide, 7.1, 4.1, 5.8, 2.95, fill_color=RGBColor(0x14, 0x2C, 0x1A))
    add_shape(slide, 7.1, 4.1, 0.06, 2.95, fill_color=C_GREEN)
    add_text_box(slide, "🔧 Specific Recommendations",
                 7.2, 4.15, 5.6, 0.4, font_size=11, bold=True, color=C_GREEN)
    rec_lines = ["1. Rate-limit: ≤50 SYN/sec per source IP toward 10.0.1.45",
                 "2. sysctl -w net.ipv4.tcp_syncookies=1",
                 "3. sysctl -w net.ipv4.tcp_synack_retries=2",
                 "4. Deploy rate-limiting toward: 10.0.1.45, 10.0.1.46"]
    add_bullet_text(slide, rec_lines, 7.2, 4.6, 5.6, 2.35,
                    font_size=10, color=C_LIGHT, bullet_char="", spacing_before=5)


def slide_usage(prs):
    """Slide 8 — How to Use"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)
    add_divider_line(slide, 0.5, 1.45, 12.33)

    add_text_box(slide, "How to Use", 0.5, 0.25, 12.0, 0.8,
                 font_size=34, bold=True, color=C_WHITE)
    add_text_box(slide, "Three interfaces — GUI, CLI, REST API — all producing the same rich HTML reports",
                 0.5, 1.05, 12.0, 0.5, font_size=16, color=C_SUBTLE)

    # Step flow
    steps = [
        ("1", "Install tshark",
         "sudo apt install tshark\n# or brew install wireshark (macOS)"),
        ("2", "Launch app",
         "./dist/AI-Wireshark-Analyzer/AI-Wireshark-Analyzer\n# or: python -m app.main"),
        ("3", "Select PCAP",
         "Browse → open .pcap or .pcapng file\n(drag & drop supported)"),
        ("4", "Choose panel",
         "WLAN · TCP/UDP · IPv6 · Protocol\nChannel Monitor · Anomaly"),
        ("5", "Apply filters",
         "IP address filter · Port filter\nMAC filter (WLAN) · Display filter"),
        ("6", "View results",
         "HTML report opens in browser\nJSON saved to results/"),
    ]

    step_w = 2.05
    for i, (num, title, detail) in enumerate(steps):
        left = 0.3 + i * 2.15
        top = 1.65
        add_shape(slide, left, top, step_w, 0.55, fill_color=C_ACCENT)
        add_text_box(slide, f"  {num}  {title}", left, top + 0.08, step_w, 0.42,
                     font_size=13, bold=True, color=C_WHITE)
        add_shape(slide, left, top + 0.55, step_w, 1.6, fill_color=C_BLUE)
        add_text_box(slide, detail, left + 0.1, top + 0.65, step_w - 0.15, 1.45,
                     font_size=10, color=C_LIGHT, wrap=True)

    # CLI examples
    add_text_box(slide, "CLI Usage Examples", 0.4, 3.55, 12.0, 0.45,
                 font_size=16, bold=True, color=C_ACCENT)

    cli_examples = [
        ("WLAN Analysis", "python3 scripts/run_wlan_analysis.py capture.pcapng [MAC]"),
        ("Single Client", "python3 scripts/run_wlan_analysis.py capture.pcapng aa:bb:cc:dd:ee:ff"),
        ("TCP/UDP", "python3 scripts/analyze_tcp_udp.py -i capture.pcap --ip 192.168.1.1"),
        ("IPv6 Deep-Dive", "python3 scripts/run_ipv6_analysis.py capture.pcapng fd12:3456::1"),
        ("Channel Monitor", "python3 scripts/run_channel_monitor.py --pcap capture.pcapng --channel 6"),
        ("Protocol CLI", "ai-wireshark analyze -i capture.pcap --protocol tcp --filter tcp.port==443"),
    ]

    for i, (label, cmd) in enumerate(cli_examples):
        row = 4.1 + (i % 3) * 1.05
        col = 0.4 if i < 3 else 6.7
        add_shape(slide, col, row, 6.0, 0.92, fill_color=RGBColor(0x0A, 0x14, 0x22))
        add_text_box(slide, label, col + 0.1, row + 0.06, 5.7, 0.3,
                     font_size=10, bold=True, color=C_SUBTLE)
        add_text_box(slide, cmd, col + 0.1, row + 0.38, 5.7, 0.5,
                     font_size=10, color=C_GREEN)


def slide_advantages(prs):
    """Slide 9 — Advantages"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)
    add_divider_line(slide, 0.5, 1.45, 12.33)

    add_text_box(slide, "Advantages", 0.5, 0.25, 12.0, 0.8,
                 font_size=34, bold=True, color=C_WHITE)
    add_text_box(slide, "Why AI-Wireshark Analyzer over manual Wireshark or cloud tools",
                 0.5, 1.05, 12.0, 0.5, font_size=16, color=C_SUBTLE)

    advantages = [
        ("🚀  Speed",
         C_ACCENT,
         ["Seconds to root cause vs. hours of manual packet inspection",
          "tshark-native bulk extraction — handles 100K+ frame captures",
          "All analysis runs locally with no upload/download latency",
          "Parallel protocol analysis (TCP + UDP + DNS simultaneously)"]),
        ("🎯  Accuracy",
         C_GREEN,
         ["Frame-sequence evidence — not just heuristics",
          "IEEE 802.11-2020 exact status/reason code mapping",
          "EAPOL Msg1-4 stall analysis to distinguish wrong PSK vs. GTK timeout",
          "CCMP Packet Number forensics for stale-PTK deadlock detection"]),
        ("🔒  Privacy & Compliance",
         C_ORANGE,
         ["100% local — captures contain sensitive data; zero cloud risk",
          "No vendor lock-in, no subscription, MIT license",
          "Suitable for enterprise, healthcare, financial sector captures",
          "No analytics or telemetry of any kind"]),
        ("📚  Expertise Embedding",
         RGBColor(0xA8, 0x55, 0xF7),
         ["Encoded IEEE 802.11, RFC 4861, BCP38 knowledge in recommendations",
          "25+ threat types with expert-level remediation steps",
          "Captures institutional knowledge — consistent across all analysts",
          "Reduces dependency on specialist availability"]),
    ]

    left_positions = [0.4, 3.55, 6.7, 9.85]
    for i, (title, color, bullets) in enumerate(advantages):
        left = left_positions[i]
        add_shape(slide, left, 1.65, 3.0, 5.55, fill_color=C_BLUE)
        add_shape(slide, left, 1.65, 3.0, 0.07, fill_color=color)
        add_text_box(slide, title, left + 0.1, 1.75, 2.8, 0.45,
                     font_size=13, bold=True, color=color)
        add_bullet_text(slide, bullets, left + 0.1, 2.28, 2.8, 4.8,
                        font_size=11, color=C_LIGHT, bullet_char="✓ ", spacing_before=5)


def slide_comparison(prs):
    """Slide 10 — Comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)
    add_divider_line(slide, 0.5, 1.45, 12.33)

    add_text_box(slide, "Tool Comparison", 0.5, 0.25, 12.0, 0.8,
                 font_size=34, bold=True, color=C_WHITE)
    add_text_box(slide, "AI-Wireshark Analyzer vs. existing approaches",
                 0.5, 1.05, 12.0, 0.5, font_size=16, color=C_SUBTLE)

    # Table header
    headers = ["Feature", "Manual Wireshark", "Cloud Tools", "AI-Wireshark Analyzer"]
    col_widths = [3.5, 2.5, 2.5, 3.0]
    col_starts = [0.3, 3.85, 6.4, 8.95]
    header_bg = [C_BLUE, C_BLUE, C_BLUE, C_ACCENT]

    for i, (hdr, w, left, bg) in enumerate(zip(headers, col_widths, col_starts, header_bg)):
        add_shape(slide, left, 1.65, w - 0.05, 0.55, fill_color=bg)
        add_text_box(slide, hdr, left + 0.08, 1.68, w - 0.15, 0.48,
                     font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    rows_data = [
        ("Root cause analysis",          "❌ Manual",       "⚠️  Generic",    "✅ Per-capture IPs"),
        ("IEEE spec references",         "❌ None",          "❌ None",          "✅ 802.11-2020 codes"),
        ("Privacy — no data upload",     "✅ Local",         "❌ Cloud required", "✅ Fully local"),
        ("Multi-protocol (WLAN+TCP+DNS)","⚠️  Separate tools","✅ Some",          "✅ Unified GUI"),
        ("Automated remediation steps",  "❌ Manual",        "⚠️  Generic",     "✅ Specific commands"),
        ("Cost",                          "Free (skill cost)","Subscription",    "✅ MIT / Free"),
        ("Works offline / air-gapped",   "✅ Yes",           "❌ No",            "✅ Yes"),
        ("HTML + JSON reports",          "❌ Manual export", "✅ Yes",           "✅ Yes"),
    ]

    row_colors = [RGBColor(0x12, 0x22, 0x33), RGBColor(0x1A, 0x2C, 0x3E)]
    for r, row in enumerate(rows_data):
        bg = row_colors[r % 2]
        for c, (cell, left, w) in enumerate(zip(row, col_starts, col_widths)):
            add_shape(slide, left, 2.28 + r * 0.62, w - 0.05, 0.58, fill_color=bg)
            cell_color = C_GREEN if "✅" in cell else C_RED if "❌" in cell else C_ORANGE if "⚠️" in cell else C_LIGHT
            add_text_box(slide, cell, left + 0.08, 2.32 + r * 0.62, w - 0.15, 0.52,
                         font_size=11, color=cell_color,
                         align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)


def slide_stats(prs):
    """Slide 11 — Key Metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)
    add_divider_line(slide, 0.5, 1.45, 12.33)

    add_text_box(slide, "Key Capabilities at a Glance", 0.5, 0.25, 12.0, 0.8,
                 font_size=34, bold=True, color=C_WHITE)

    stats = [
        ("40+", "IEEE 802.11\nstatus/reason codes", C_ACCENT),
        ("25+", "Threat types with\ndynamic RCA", C_GREEN),
        ("9", "Analysis panels\nin GUI", C_ORANGE),
        ("100%", "Local processing\nno cloud", RGBColor(0xA8, 0x55, 0xF7)),
        ("1.1 GB", "Self-contained\nLinux binary", C_LIGHT),
        ("5", "Output formats\nHTML/JSON/PDF/CSV/PPTX", C_ACCENT),
    ]

    for i, (val, label, color) in enumerate(stats):
        left = 0.5 + (i % 3) * 4.1
        top = 1.65 if i < 3 else 3.8
        add_stat_box(slide, left, top, val, label, color=color)

    # Bottom highlight bar
    add_shape(slide, 0.4, 6.1, 12.53, 1.1, fill_color=RGBColor(0x1E, 0x40, 0xAF))
    add_shape(slide, 0.4, 6.1, 0.08, 1.1, fill_color=C_ACCENT)

    highlights = [
        "IEEE 802.11-2020 §12.4 (WPA3-SAE)",
        "RFC 2827 / BCP38 (IP spoofing prevention)",
        "802.11w MFP / PMF (management frame protection)",
        "RFC 6105 RA Guard / IPv6 FHS",
        "EAPOL RFC 8023 / IEEE 802.1X",
    ]
    add_text_box(slide, "Standards Referenced in Recommendations:",
                 0.6, 6.15, 3.5, 0.45, font_size=11, bold=True, color=C_ACCENT)
    add_text_box(slide, "  ·  ".join(highlights),
                 4.2, 6.25, 8.5, 0.75, font_size=11, color=C_LIGHT)


def slide_roadmap(prs):
    """Slide 12 — Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)
    add_divider_line(slide, 0.5, 1.45, 12.33)

    add_text_box(slide, "Roadmap", 0.5, 0.25, 12.0, 0.8,
                 font_size=34, bold=True, color=C_WHITE)
    add_text_box(slide, "Planned improvements in priority order",
                 0.5, 1.05, 12.0, 0.5, font_size=16, color=C_SUBTLE)

    phases = [
        ("v1.8.0  —  Q3 2026", C_GREEN,
         [
             "Windows x64 binary distribution (native build host)",
             "macOS app bundle / DMG distribution",
             "Live capture support: tshark -i wlan0 piped directly into all panels",
             "PDF export for HTML reports (for incident documentation)",
             "Enhanced TCP/UDP timeline charts with per-flow breakdown",
         ]),
        ("v1.9.0  —  Q4 2026", C_ACCENT,
         [
             "Multi-capture comparison: diff two captures side-by-side",
             "Automated baseline profiling — flag deviations from normal traffic",
             "REST API hardening: JWT auth, rate-limiting, OpenAPI 3.0 spec",
             "Plugin architecture: custom threat detectors as external modules",
             "Wireshark 4.x coloring rules export from analysis results",
         ]),
        ("v2.0.0  —  H1 2027", C_ORANGE,
         [
             "AI-assisted threat correlation across multiple capture sessions",
             "LLM-powered natural language query: 'show all Wi-Fi auth failures today'",
             "Vendor firmware signature database for AP bug identification",
             "PCAP anonymization engine (MAC/IP anonymization before sharing)",
             "Dashboard with historical trend analysis across capture archive",
         ]),
        ("Ongoing", C_SUBTLE,
         [
             "Expand IEEE 802.11 status/reason code coverage (new 802.11be / EHT codes)",
             "Automated test suite expansion (capture-based regression tests)",
             "Documentation: video walkthroughs per panel and threat type",
             "Community: public capture library for testing and demonstration",
             "Performance: streaming analysis for captures > 1 GB",
         ]),
    ]

    for i, (title, color, bullets) in enumerate(phases):
        left = 0.4 + (i % 2) * 6.47
        top = 1.65 + (i // 2) * 2.85
        add_card(slide, left, top, 6.2, 2.65, title, bullets,
                 card_bg=C_BLUE, title_color=color,
                 title_size=13, body_size=10.5)


def slide_importance(prs):
    """Slide 13 — Why It Matters"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)
    add_divider_line(slide, 0.5, 1.45, 12.33)

    add_text_box(slide, "Why It Matters", 0.5, 0.25, 12.0, 0.8,
                 font_size=34, bold=True, color=C_WHITE)

    scenarios = [
        ("🏢  Enterprise IT",
         "Wi-Fi authentication failures impact productivity. A network engineer can diagnose \"WPA3-SAE PMKID stale cache\" vs \"wrong PSK\" in seconds instead of hours, pinpoint the exact AP and client, and get a corrective command — without specialist WPA3 expertise."),
        ("🏥  Healthcare / Critical Infrastructure",
         "Medical devices on Wi-Fi cannot send telemetry to cloud tools. AI-Wireshark Analyzer runs air-gapped, never uploads PCAP data, and identifies DHCP starvation attacks or rogue DHCP servers that could disconnect life-critical monitoring equipment."),
        ("🔐  Security Teams / SOC",
         "DNS tunneling, SYN floods, and port scans are detected with specific attacker IPs ready to add to firewall block lists. No cross-referencing of multiple tools — one report with targeted iptables/nftables commands."),
        ("🎓  Training & Education",
         "Each finding references the IEEE spec section responsible. Analysts learn protocol fundamentals from real captures, with the tool acting as an interactive guide — a textbook that reads your specific capture."),
    ]

    for i, (title, body) in enumerate(scenarios):
        left = 0.4 + (i % 2) * 6.47
        top = 1.65 + (i // 2) * 2.75

        add_shape(slide, left, top, 6.2, 2.55, fill_color=C_BLUE)
        add_shape(slide, left, top, 0.06, 2.55, fill_color=C_ACCENT)
        add_text_box(slide, title, left + 0.15, top + 0.1, 5.9, 0.45,
                     font_size=14, bold=True, color=C_ACCENT)
        add_text_box(slide, body, left + 0.15, top + 0.6, 5.9, 1.85,
                     font_size=11, color=C_LIGHT, wrap=True)


def slide_closing(prs):
    """Slide 14 — Closing / Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, C_NAVY)

    # Full-width top accent bar
    add_shape(slide, 0, 0, 13.33, 0.1, fill_color=C_ACCENT)
    add_shape(slide, 0, 7.4, 13.33, 0.1, fill_color=C_ACCENT)

    # Large icon
    add_text_box(slide, "📡", 5.6, 0.8, 2.2, 1.5,
                 font_size=72, color=C_ACCENT, align=PP_ALIGN.CENTER)

    add_text_box(slide, "AI-Wireshark Analyzer", 1.0, 2.3, 11.33, 1.0,
                 font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Intelligent · Local · Actionable",
                 1.0, 3.25, 11.33, 0.65,
                 font_size=22, color=C_ACCENT, align=PP_ALIGN.CENTER)

    add_text_box(slide, "v1.7.1  ·  MIT License  ·  Linux x64 build available",
                 1.0, 4.0, 11.33, 0.5,
                 font_size=14, color=C_SUBTLE, align=PP_ALIGN.CENTER)

    # CTA boxes
    ctas = [
        ("🚀 Get Started", "Download the Linux binary\n(1.1 GB, no install needed)", C_ACCENT),
        ("📋 Run Analysis", "Open any PCAP → choose panel\n→ instant HTML report", C_GREEN),
        ("🗺️  Contribute", "Follow the roadmap or add\nyour own protocol analyzer", C_ORANGE),
    ]
    for i, (title, body, color) in enumerate(ctas):
        left = 1.0 + i * 3.9
        add_shape(slide, left, 4.75, 3.55, 1.7, fill_color=C_BLUE)
        add_shape(slide, left, 4.75, 3.55, 0.07, fill_color=color)
        add_text_box(slide, title, left + 0.15, 4.85, 3.2, 0.45,
                     font_size=13, bold=True, color=color)
        add_text_box(slide, body, left + 0.15, 5.35, 3.2, 1.0,
                     font_size=11, color=C_LIGHT, wrap=True)

    add_text_box(slide, "Requirements: tshark installed · Linux/Windows/macOS desktop session",
                 1.0, 6.8, 11.33, 0.45,
                 font_size=11, color=C_SUBTLE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def build_presentation(output_path: str) -> str:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_title(prs);       print("  [1/14] Title")
    slide_problem(prs);     print("  [2/14] Problem Statement")
    slide_solution(prs);    print("  [3/14] Solution Overview")
    slide_architecture(prs);print("  [4/14] Architecture")
    slide_wlan(prs);        print("  [5/14] WLAN Analysis")
    slide_protocol(prs);    print("  [6/14] Protocol Analyzers")
    slide_reporting(prs);   print("  [7/14] Intelligent Reporting")
    slide_usage(prs);       print("  [8/14] Usage / How To Use")
    slide_advantages(prs);  print("  [9/14] Advantages")
    slide_comparison(prs);  print(" [10/14] Tool Comparison")
    slide_stats(prs);       print(" [11/14] Key Metrics")
    slide_roadmap(prs);     print(" [12/14] Roadmap")
    slide_importance(prs);  print(" [13/14] Why It Matters")
    slide_closing(prs);     print(" [14/14] Closing / CTA")

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    output = sys.argv[1] if len(sys.argv) > 1 else "docs/AI-Wireshark-Analyzer-Presentation.pptx"
    result = build_presentation(output)
    print(f"\nPresentation saved: {result}")
    print(f"File size: {Path(result).stat().st_size / 1024:.0f} KB")
    print(f"Slides: 14")
